#!/usr/bin/env python3
"""Audit presentation readability and empty PowerPoint placeholders.

Usage:
  python3 skills/presentation-redesign/scripts/audit_pptx.py deck.pptx
  python3 skills/presentation-redesign/scripts/audit_pptx.py deck.pptx \
      --remove-empty-placeholders --output cleaned.pptx

The command prints JSON and exits 1 when visible text falls below the requested
minimum or an empty placeholder remains. Theme-inherited font sizes are treated
as failures because their rendered size is not provable from slide content.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def inspect_presentation(path: Path, min_font_size: float, remove_empty: bool):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("Missing dependency: pip install python-pptx") from exc

    presentation = Presentation(path)
    empty_placeholders: list[dict[str, object]] = []
    small_text: list[dict[str, object]] = []
    removed = 0

    def inspect_text_frame(slide_number: int, source: str, text_frame) -> None:
        for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
            for run in paragraph.runs:
                value = run.text.strip()
                if not value:
                    continue
                size = run.font.size.pt if run.font.size is not None else None
                if size is None or size < min_font_size:
                    small_text.append(
                        {
                            "slide": slide_number,
                            "source": source,
                            "paragraph": paragraph_index,
                            "size_pt": size,
                            "text": value,
                        }
                    )

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.has_text_frame and not shape.text.strip():
                finding = {
                    "slide": slide_number,
                    "shape": shape.name,
                    "placeholder_type": str(shape.placeholder_format.type),
                }
                if remove_empty:
                    shape._element.getparent().remove(shape._element)
                    removed += 1
                    continue
                empty_placeholders.append(finding)

            if shape.has_text_frame:
                inspect_text_frame(slide_number, shape.name, shape.text_frame)
            if shape.has_table:
                for row_index, row in enumerate(shape.table.rows):
                    for cell_index, cell in enumerate(row.cells):
                        inspect_text_frame(slide_number, f"table[{row_index},{cell_index}]", cell.text_frame)

    return presentation, {
        "deck": str(path),
        "min_font_size": min_font_size,
        "empty_placeholders": empty_placeholders,
        "small_text": small_text,
        "removed_empty_placeholders": removed,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck", type=Path)
    parser.add_argument("--min-font-size", type=float, default=28.0)
    parser.add_argument("--remove-empty-placeholders", action="store_true")
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    if args.remove_empty_placeholders and args.output is None:
        parser.error("--output is required with --remove-empty-placeholders")

    presentation, report = inspect_presentation(
        args.deck, args.min_font_size, args.remove_empty_placeholders
    )
    if args.output is not None:
        presentation.save(args.output)
        report["output"] = str(args.output)

    print(json.dumps(report, indent=2))
    return 1 if report["empty_placeholders"] or report["small_text"] else 0


if __name__ == "__main__":
    sys.exit(main())
